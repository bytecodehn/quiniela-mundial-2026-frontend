"""
Export quiniela-mundial-2026.html → PPTX with footer-rail discipline.
Canvas: 13.333" x 7.5" (16:9)
Rail constants:
  CONTENT_MAX_Y = 6.70"  — no content element below this line
  FOOTER_TOP    = 6.85"  — footer rail starts here
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import copy

# ── Rail constants (inches → EMU) ──────────────────────────────────────────
INCH = 914400
CONTENT_MAX_Y = int(6.70 * INCH)
FOOTER_TOP    = int(6.85 * INCH)
SLIDE_W       = int(13.333 * INCH)
SLIDE_H       = int(7.5 * INCH)

# ── Design system colors (from tokens.css, approximated to sRGB) ──────────
MIDNIGHT      = RGBColor(0x0C, 0x11, 0x24)  # --bg-primary oklch(12% 0.02 260)
DARK_SURFACE  = RGBColor(0x16, 0x1D, 0x33)  # --bg-surface oklch(18% 0.018 260)
ELEVATED      = RGBColor(0x1C, 0x24, 0x3A)  # --bg-elevated oklch(22% 0.018 260)
EMERALD       = RGBColor(0x2E, 0xC5, 0x6D)  # --green oklch(60% 0.18 145)
GOLD          = RGBColor(0xF0, 0xC6, 0x3B)  # --gold oklch(75% 0.15 85)
WHITE         = RGBColor(0xF0, 0xF1, 0xF3)  # --fg oklch(93% 0.005 260)
MUTED         = RGBColor(0x68, 0x70, 0x83)  # --fg-muted oklch(45% 0.015 260)
SECONDARY     = RGBColor(0x9C, 0xA1, 0xB0)  # --fg-secondary oklch(65% 0.015 260)
RED           = RGBColor(0xE5, 0x4C, 0x4C)  # --red
BLUE          = RGBColor(0x4A, 0x8F, 0xE8)  # --blue
CYAN          = RGBColor(0x54, 0xC7, 0xD6)  # --cyan
SHELL_BG      = RGBColor(0x08, 0x09, 0x0D)  # shell background

# ── Font constants ──────────────────────────────────────────────────────────
LATIN_FACE = "SF Pro Display"
EA_FACE    = "Yu Gothic UI"  # explicit CJK fallback — never Microsoft JhengHei
MONO_FACE  = "SF Mono"

# ── Helper: set latin + ea typefaces on a run ──────────────────────────────
def set_run_typefaces(run, latin=LATIN_FACE, ea=EA_FACE, italic=False, bold=False, size=None, color=None):
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:latin'), latin) if latin else None
    rPr.set(qn('a:ea'), ea) if ea else None
    if italic:
        rPr.set('i', '1')
    if bold:
        rPr.set('b', '1')
    if size:
        rPr.set('sz', str(int(size * 100)))  # hundredths of a point
    if color:
        srgb = rPr.makeelement(qn('a:srgbClr'), {'val': color})
        rPr.append(srgb)

def set_para_typefaces(para, latin=LATIN_FACE, ea=EA_FACE):
    """Set typeface on all runs in a paragraph."""
    for run in para.runs:
        set_run_typefaces(run, latin=latin, ea=ea)


# ── Cursor: y-position tracker that enforces the footer rail ──────────────
class Cursor:
    def __init__(self, start_y=Emu(0)):
        self.y = start_y
        self.max_y = CONTENT_MAX_Y

    def advance(self, dy_emu):
        next_y = self.y + dy_emu
        if next_y > self.max_y:
            raise ValueError(
                f"Cursor would cross CONTENT_MAX_Y ({self.max_y/INCH:.2f}in): "
                f"{self.y/INCH:.2f} + {dy_emu/INCH:.2f} = {next_y/INCH:.2f} > {self.max_y/INCH:.2f}"
            )
        self.y = next_y

    def try_advance(self, dy_emu):
        next_y = self.y + dy_emu
        if next_y > self.max_y:
            return False
        self.y = next_y
        return True

    def remaining(self):
        return self.max_y - self.y

    def reset(self, y=Emu(0)):
        self.y = y


# ── Shape helpers ──────────────────────────────────────────────────────────
def add_bg(slide, color):
    """Fill entire slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_face=None, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    if font_face:
        set_run_typefaces(run, latin=font_face, ea=EA_FACE)
    else:
        set_run_typefaces(run)
    return txBox

def add_rich_paragraph(text_frame, runs_data, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    """Add a paragraph with multiple styled runs. runs_data: list of (text, kwargs)."""
    if text_frame.paragraphs and text_frame.paragraphs[0].text == '':
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    for text, kwargs in runs_data:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(kwargs.get('size', 18))
        run.font.color.rgb = kwargs.get('color', WHITE)
        run.font.bold = kwargs.get('bold', False)
        run.font.italic = kwargs.get('italic', False)
        latin = kwargs.get('latin', LATIN_FACE)
        ea = kwargs.get('ea', EA_FACE)
        set_run_typefaces(run, latin=latin, ea=ea,
                          italic=run.font.italic, bold=run.font.bold,
                          size=kwargs.get('size'), color=None)
    return p

def add_multiline(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, kwargs_dict) tuples, each becomes a paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    for i, (text, kwargs) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = kwargs.get('alignment', PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(kwargs.get('size', 18))
        run.font.color.rgb = kwargs.get('color', WHITE)
        run.font.bold = kwargs.get('bold', False)
        run.font.italic = kwargs.get('italic', False)
        set_run_typefaces(run, latin=kwargs.get('latin', LATIN_FACE), ea=kwargs.get('ea', EA_FACE),
                          italic=run.font.italic, bold=run.font.bold, size=kwargs.get('size'))
    return txBox


# ── Slide builders ─────────────────────────────────────────────────────────
def build_slide_1_cover(prs):
    """Cover: Quiniela Mundial 2026"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, MIDNIGHT)

    # Accent bar at top
    add_rect(slide, 0, 0, SLIDE_W, int(0.06 * INCH), fill_color=EMERALD)

    # Center content vertically using budget centering (not MARGIN_TOP)
    # The content zone is 0 to CONTENT_MAX_Y; we'll place the main block centered
    content_block_h = int(3.5 * INCH)
    block_top = int((CONTENT_MAX_Y - content_block_h) / 2)

    # Trophy icon area
    add_text_box(slide, 0, block_top, SLIDE_W, int(0.8 * INCH),
                 "🏆", font_size=48, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, int(1.5*INCH), block_top + int(0.9*INCH),
                 SLIDE_W - int(3*INCH), int(1.2*INCH),
                 "Quiniela Mundial 2026", font_size=44, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, int(2*INCH), block_top + int(2.2*INCH),
                 SLIDE_W - int(4*INCH), int(0.7*INCH),
                 "Pronostica · Compite · Gana", font_size=22,
                 color=EMERALD, alignment=PP_ALIGN.CENTER)

    # Tagline
    add_text_box(slide, int(2.5*INCH), block_top + int(3.0*INCH),
                 SLIDE_W - int(5*INCH), int(0.6*INCH),
                 "Plataforma de predicciones deportivas — Mundial 2026",
                 font_size=16, color=SECONDARY, alignment=PP_ALIGN.CENTER)

    # Footer rail
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, int(0.65 * INCH),
             fill_color=DARK_SURFACE)
    add_text_box(slide, int(1*INCH), FOOTER_TOP + int(0.12*INCH),
                 SLIDE_W - int(2*INCH), int(0.4*INCH),
                 "Sin apuestas · Sin dinero real · Solo diversión competitiva",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build_slide_2_design_system(prs):
    """Design System overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MIDNIGHT)
    cursor = Cursor(start_y=int(0.6 * INCH))

    # Slide label
    add_rect(slide, 0, 0, int(0.06*INCH), int(0.06*INCH), fill_color=EMERALD)
    add_text_box(slide, int(0.8*INCH), int(0.15*INCH), int(6*INCH), int(0.5*INCH),
                 "Design System", font_size=14, color=MUTED, bold=True)

    # Title
    add_text_box(slide, int(0.8*INCH), cursor.y, int(8*INCH), int(0.8*INCH),
                 "Sistema de diseño · Quiniela Mundial 2026", font_size=36, bold=True)
    cursor.advance(int(0.9 * INCH))

    # Palette row
    palette_y = cursor.y
    swatch_size = int(0.7 * INCH)
    gap = int(0.2 * INCH)
    start_x = int(0.8 * INCH)
    colors = [
        ("Midnight", MIDNIGHT),
        ("Surface", DARK_SURFACE),
        ("Elevated", ELEVATED),
        ("Emerald", EMERALD),
        ("Gold", GOLD),
        ("White", WHITE),
        ("Muted", MUTED),
        ("Red", RED),
        ("Blue", BLUE),
    ]
    for i, (name, c) in enumerate(colors):
        x = start_x + i * (swatch_size + gap)
        add_rect(slide, x, palette_y, swatch_size, swatch_size, fill_color=c)
        add_rect(slide, x, palette_y, swatch_size, swatch_size,
                 border_color=RGBColor(0xFF, 0xFF, 0xFF), border_width=Pt(0.5))
        add_text_box(slide, x, palette_y + swatch_size + int(0.05*INCH),
                     swatch_size, int(0.3*INCH),
                     name, font_size=9, color=SECONDARY, alignment=PP_ALIGN.CENTER)

    cursor.advance(int(1.2 * INCH))

    # Typography section
    add_text_box(slide, int(0.8*INCH), cursor.y, int(10*INCH), int(0.4*INCH),
                 "Tipografía", font_size=20, bold=True, color=WHITE)
    cursor.advance(int(0.5 * INCH))

    font_info_y = cursor.y
    fonts = [
        ("Display", "SF Pro Display, system-ui, sans-serif", WHITE),
        ("Body", "SF Pro Text, system-ui, sans-serif", SECONDARY),
        ("Mono", "SF Mono, JetBrains Mono, monospace", CYAN),
        ("CJK", "Yu Gothic UI (a:ea explicit)", MUTED),
    ]
    for i, (label, info, clr) in enumerate(fonts):
        y = font_info_y + i * int(0.35 * INCH)
        add_text_box(slide, int(0.8*INCH), y, int(2.5*INCH), int(0.3*INCH),
                     label, font_size=13, bold=True, color=clr)
        add_text_box(slide, int(3.5*INCH), y, int(6*INCH), int(0.3*INCH),
                     info, font_size=13, color=SECONDARY)

    cursor.advance(int(1.6 * INCH))

    # Key posture notes
    add_text_box(slide, int(0.8*INCH), cursor.y, int(10*INCH), int(0.4*INCH),
                 "Postura visual", font_size=20, bold=True, color=WHITE)
    cursor.advance(int(0.5 * INCH))

    postures = [
        "Fondo midnight blue · Cards con bordes sutiles · Sin sombras grandes",
        "Un acento verde esmeralda · Usado máximo 2 veces por pantalla",
        "Diseño responsive: desktop sidebar + mobile bottom nav",
        "14 pantallas · Arquitectura lista para Next.js + TypeScript",
    ]
    for i, ptext in enumerate(postures):
        y = cursor.y + i * int(0.3 * INCH)
        add_text_box(slide, int(0.8*INCH), y, int(10*INCH), int(0.3*INCH),
                     f"·  {ptext}", font_size=13, color=SECONDARY)

    # Footer rail
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, int(0.65 * INCH),
             fill_color=DARK_SURFACE)
    add_text_box(slide, int(1*INCH), FOOTER_TOP + int(0.12*INCH),
                 SLIDE_W - int(2*INCH), int(0.4*INCH),
                 "02 / 08  ·  Sistema de diseño",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build_slide_3_public_screens(prs):
    """Public screens: Landing, Login, Register"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MIDNIGHT)
    cursor = Cursor(start_y=int(0.6 * INCH))

    add_rect(slide, 0, 0, int(0.06*INCH), int(0.06*INCH), fill_color=EMERALD)
    add_text_box(slide, int(0.8*INCH), int(0.15*INCH), int(6*INCH), int(0.5*INCH),
                 "Pantallas públicas", font_size=14, color=MUTED, bold=True)
    add_text_box(slide, int(0.8*INCH), cursor.y, int(8*INCH), int(0.8*INCH),
                 "Landing · Inicio de sesión · Registro", font_size=36, bold=True)
    cursor.advance(int(1.0 * INCH))

    # Three cards overview
    cards_data = [
        ("🏠", "Landing Page",
         "Hero con propuesta de valor · CTA principal\nSección «Cómo funciona» · Beneficios\nVista previa del ranking global\nPróximos partidos destacados\nFooter con navegación completa"),
        ("🔐", "Iniciar sesión",
         "Formulario de email + contraseña\nRecuperar contraseña\nGoogle OAuth\nDividor visual con «O continúa con»\nDiseño centrado, limpio y moderno"),
        ("📝", "Registro",
         "Formulario: nombre, email, país, contraseña\nConfirmación de contraseña\nTérminos y condiciones\nPaís / selección favorita\nConsistente visualmente con login"),
    ]

    card_w = int(3.75 * INCH)
    card_h = int(3.8 * INCH)
    card_gap = int(0.3 * INCH)
    total_w = 3 * card_w + 2 * card_gap
    start_x = int((SLIDE_W - total_w) / 2)

    for i, (icon, title, desc) in enumerate(cards_data):
        x = start_x + i * (card_w + card_gap)
        y = cursor.y
        # Card bg
        add_rect(slide, x, y, card_w, card_h, fill_color=DARK_SURFACE,
                 border_color=RGBColor(0x28, 0x30, 0x48))
        # Icon
        add_text_box(slide, x + int(0.2*INCH), y + int(0.2*INCH),
                     card_w - int(0.4*INCH), int(0.6*INCH),
                     icon, font_size=32, alignment=PP_ALIGN.LEFT)
        # Title
        add_text_box(slide, x + int(0.2*INCH), y + int(0.8*INCH),
                     card_w - int(0.4*INCH), int(0.4*INCH),
                     title, font_size=18, bold=True, color=WHITE)
        # Description
        add_text_box(slide, x + int(0.2*INCH), y + int(1.3*INCH),
                     card_w - int(0.4*INCH), int(2.3*INCH),
                     desc, font_size=11, color=SECONDARY)
        # Badge
        badge_y = y + card_h - int(0.4*INCH)
        add_rect(slide, x + int(0.2*INCH), badge_y, int(1.0*INCH), int(0.26*INCH),
                 fill_color=RGBColor(0x1A, 0x28, 0x3A))
        add_text_box(slide, x + int(0.2*INCH), badge_y + int(0.03*INCH),
                     int(1.0*INCH), int(0.2*INCH),
                     "PÚBLICO", font_size=8, color=BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Footer rail
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, int(0.65 * INCH),
             fill_color=DARK_SURFACE)
    add_text_box(slide, int(1*INCH), FOOTER_TOP + int(0.12*INCH),
                 SLIDE_W - int(2*INCH), int(0.4*INCH),
                 "03 / 08  ·  Pantallas públicas",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build_slide_4_dashboard(prs):
    """Dashboard overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MIDNIGHT)
    cursor = Cursor(start_y=int(0.6 * INCH))

    add_rect(slide, 0, 0, int(0.06*INCH), int(0.06*INCH), fill_color=EMERALD)
    add_text_box(slide, int(0.8*INCH), int(0.15*INCH), int(6*INCH), int(0.5*INCH),
                 "Usuario autenticado", font_size=14, color=MUTED, bold=True)
    add_text_box(slide, int(0.8*INCH), cursor.y, int(10*INCH), int(0.8*INCH),
                 "Dashboard principal", font_size=36, bold=True)
    cursor.advance(int(0.9 * INCH))

    # Left column — dashboard features
    features = [
        ("📊", "Resumen de puntos", "Puntuación total, posición global y en tu grupo"),
        ("⚽", "Partidos próximos", "Próximos partidos con botón para predecir"),
        ("📋", "Predicciones recientes", "Últimos resultados, aciertos, puntos obtenidos"),
        ("🏆", "Ranking", "Posición destacada con tendencia"),
    ]
    col_w = int(5.5 * INCH)
    for i, (icon, title, desc) in enumerate(features):
        y = cursor.y + i * int(0.65 * INCH)
        add_text_box(slide, int(0.8*INCH), y, int(0.5*INCH), int(0.5*INCH),
                     icon, font_size=24)
        add_text_box(slide, int(1.5*INCH), y, col_w, int(0.3*INCH),
                     title, font_size=16, bold=True, color=WHITE)
        add_text_box(slide, int(1.5*INCH), y + int(0.28*INCH), col_w, int(0.3*INCH),
                     desc, font_size=12, color=SECONDARY)

    cursor.advance(int(3.0 * INCH))

    # Quick actions section
    add_text_box(slide, int(0.8*INCH), cursor.y, int(8*INCH), int(0.4*INCH),
                 "Acceso rápido desde el dashboard:", font_size=16, bold=True, color=EMERALD)
    cursor.advance(int(0.45 * INCH))

    actions = "Hacer predicción  ·  Ver leaderboard  ·  Grupos  ·  Calendario de partidos"
    add_text_box(slide, int(0.8*INCH), cursor.y, int(10*INCH), int(0.4*INCH),
                 actions, font_size=14, color=CYAN)

    # Right side — mockup sketch
    mockup_x = int(7.5 * INCH)
    mockup_y = int(1.5 * INCH)
    mockup_w = int(5.2 * INCH)
    mockup_h = int(3.8 * INCH)
    add_rect(slide, mockup_x, mockup_y, mockup_w, mockup_h,
             fill_color=DARK_SURFACE, border_color=RGBColor(0x28, 0x30, 0x48))

    # Mockup header
    add_rect(slide, mockup_x, mockup_y, mockup_w, int(0.45*INCH),
             fill_color=ELEVATED)
    add_text_box(slide, mockup_x + int(0.15*INCH), mockup_y + int(0.08*INCH),
                 int(2*INCH), int(0.3*INCH),
                 "Dashboard", font_size=12, bold=True, color=WHITE)
    add_text_box(slide, mockup_x + int(2*INCH), mockup_y + int(0.08*INCH),
                 int(2*INCH), int(0.3*INCH),
                 "👤 Juan Pérez", font_size=10, color=SECONDARY)

    # Mockup stat bars
    stat_positions = [(0.15, 0.6), (1.4, 0.6), (2.65, 0.6), (3.9, 0.6)]
    for sx, sy in stat_positions:
        sx_emu = mockup_x + int(sx * INCH)
        sy_emu = mockup_y + int(sy * INCH)
        add_rect(slide, sx_emu, sy_emu, int(1.1*INCH), int(0.75*INCH),
                 fill_color=ELEVATED, border_color=RGBColor(0x28, 0x30, 0x48))

    # Mockup stat labels
    stat_labels = [("1,280", "pts"), ("#24", "global"), ("#3", "grupo"), ("8/12", "aciertos")]
    for i, (val, lbl) in enumerate(stat_labels):
        sx, sy = stat_positions[i]
        sx_emu = mockup_x + int(sx * INCH) + int(0.1*INCH)
        sy_emu = mockup_y + int(sy * INCH) + int(0.08*INCH)
        add_text_box(slide, sx_emu, sy_emu, int(0.9*INCH), int(0.35*INCH),
                     val, font_size=16, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, sx_emu, sy_emu + int(0.35*INCH), int(0.9*INCH), int(0.3*INCH),
                     lbl, font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Mockup upcoming matches section
    mm_y = mockup_y + int(1.5 * INCH)
    add_text_box(slide, mockup_x + int(0.15*INCH), mm_y,
                 int(2*INCH), int(0.25*INCH),
                 "Próximos partidos", font_size=10, bold=True, color=SECONDARY)
    match_rows = [
        "🇦🇷 Argentina vs 🇧🇷 Brasil",
        "🇪🇸 España vs 🇩🇪 Alemania",
        "🇫🇷 Francia vs 🇮🇹 Italia",
    ]
    for i, row in enumerate(match_rows):
        ry = mm_y + int(0.3*INCH) + i * int(0.25*INCH)
        add_text_box(slide, mockup_x + int(0.15*INCH), ry,
                     int(4.5*INCH), int(0.25*INCH),
                     row, font_size=9, color=SECONDARY)

    # Footer rail
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, int(0.65 * INCH),
             fill_color=DARK_SURFACE)
    add_text_box(slide, int(1*INCH), FOOTER_TOP + int(0.12*INCH),
                 SLIDE_W - int(2*INCH), int(0.4*INCH),
                 "04 / 08  ·  Dashboard del usuario",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build_slide_5_matches_predictions(prs):
    """Matches & Predictions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MIDNIGHT)
    cursor = Cursor(start_y=int(0.6 * INCH))

    add_rect(slide, 0, 0, int(0.06*INCH), int(0.06*INCH), fill_color=EMERALD)
    add_text_box(slide, int(0.8*INCH), int(0.15*INCH), int(6*INCH), int(0.5*INCH),
                 "Flujo de predicciones", font_size=14, color=MUTED, bold=True)
    add_text_box(slide, int(0.8*INCH), cursor.y, int(10*INCH), int(0.8*INCH),
                 "Calendario de partidos · Pronosticar · Mis predicciones",
                 font_size=32, bold=True)
    cursor.advance(int(0.9 * INCH))

    # Three column layout
    col_w = int(3.8 * INCH)
    col_gap = int(0.25 * INCH)
    cols_x = [int(0.6*INCH), int(0.6*INCH) + col_w + col_gap, int(0.6*INCH) + 2*(col_w + col_gap)]
    col_h = int(3.6 * INCH)

    col_data = [
        ("📅", "Calendario de partidos", [
            "Listado completo con filtros por grupo",
            "Filtros por fecha y estado del partido",
            "Tarjetas con equipo local, visitante",
            "Fecha, hora, estadio",
            "Estados: próximo, cerrado, finalizado",
            "Botón «Predecir» en cada partido",
            "Fase de grupos y eliminatorias",
        ]),
        ("🎯", "Pronosticar partido", [
            "Equipos enfrentados con escudos",
            "Fecha, hora y countdown",
            "Inputs de marcador local y visitante",
            "Botón guardar predicción",
            "Estado de la predicción",
            "Mensaje si ya no se puede editar",
            "Reglas de puntuación visibles",
        ]),
        ("📋", "Mis predicciones", [
            "Historial completo de predicciones",
            "Marcador pronosticado vs real",
            "Puntos obtenidos por acierto",
            "Estados: pendiente, acertada, fallida, exacta",
            "Filtros por fase, fecha, estado",
            "Estadísticas de precisión",
            "Feed de resultados recientes",
        ]),
    ]

    for i, (icon, title, items) in enumerate(col_data):
        x = cols_x[i]
        y = cursor.y
        # Column card
        add_rect(slide, x, y, col_w, col_h, fill_color=DARK_SURFACE,
                 border_color=RGBColor(0x28, 0x30, 0x48))
        # Icon + title
        add_text_box(slide, x + int(0.15*INCH), y + int(0.15*INCH),
                     int(0.5*INCH), int(0.5*INCH),
                     icon, font_size=28)
        add_text_box(slide, x + int(0.65*INCH), y + int(0.2*INCH),
                     col_w - int(0.8*INCH), int(0.4*INCH),
                     title, font_size=17, bold=True, color=WHITE)
        # Items
        for j, item in enumerate(items):
            iy = y + int(0.7*INCH) + j * int(0.35*INCH)
            add_text_box(slide, x + int(0.2*INCH), iy, col_w - int(0.4*INCH), int(0.3*INCH),
                         f"·  {item}", font_size=11, color=SECONDARY)

    # Footer rail
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, int(0.65 * INCH),
             fill_color=DARK_SURFACE)
    add_text_box(slide, int(1*INCH), FOOTER_TOP + int(0.12*INCH),
                 SLIDE_W - int(2*INCH), int(0.4*INCH),
                 "05 / 08  ·  Calendario · Predicción · Historial",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build_slide_6_leaderboard_groups(prs):
    """Leaderboard & Groups"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MIDNIGHT)
    cursor = Cursor(start_y=int(0.6 * INCH))

    add_rect(slide, 0, 0, int(0.06*INCH), int(0.06*INCH), fill_color=EMERALD)
    add_text_box(slide, int(0.8*INCH), int(0.15*INCH), int(6*INCH), int(0.5*INCH),
                 "Social & competitivo", font_size=14, color=MUTED, bold=True)
    add_text_box(slide, int(0.8*INCH), cursor.y, int(10*INCH), int(0.8*INCH),
                 "Leaderboard · Grupos privados", font_size=36, bold=True)
    cursor.advance(int(0.9 * INCH))

    # Two panel layout
    left_w = int(5.5 * INCH)
    right_w = int(5.5 * INCH)
    gap = int(0.4 * INCH)
    left_x = int(0.6 * INCH)
    right_x = left_x + left_w + gap
    panel_h = int(3.8 * INCH)

    # Left panel — Leaderboard
    add_rect(slide, left_x, cursor.y, left_w, panel_h, fill_color=DARK_SURFACE,
             border_color=RGBColor(0x28, 0x30, 0x48))
    add_text_box(slide, left_x + int(0.2*INCH), cursor.y + int(0.15*INCH),
                 int(3*INCH), int(0.4*INCH),
                 "🏆  Ranking global", font_size=18, bold=True, color=WHITE)

    # Podium sketch
    podium_y = cursor.y + int(0.6*INCH)
    podium_items = [
        ("🥇", GOLD, "1,420 pts"),
        ("🥈", SECONDARY, "1,180 pts"),
        ("🥉", RGBColor(0xCD, 0x7F, 0x32), "1,050 pts"),
    ]
    for i, (medal, clr, pts) in enumerate(podium_items):
        px = left_x + int(0.2*INCH) + i * int(1.7*INCH)
        add_rect(slide, px, podium_y, int(1.5*INCH), int(1.3*INCH),
                 fill_color=ELEVATED, border_color=RGBColor(0x28, 0x30, 0x48))
        add_text_box(slide, px, podium_y + int(0.1*INCH),
                     int(1.5*INCH), int(0.4*INCH),
                     medal, font_size=24, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, px, podium_y + int(0.5*INCH),
                     int(1.5*INCH), int(0.3*INCH),
                     f"Jugador {i+1}", font_size=12, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, px, podium_y + int(0.85*INCH),
                     int(1.5*INCH), int(0.3*INCH),
                     pts, font_size=10, color=clr,
                     alignment=PP_ALIGN.CENTER)

    # Rank table sketch
    table_y = cursor.y + int(2.1*INCH)
    headers = ["#", "Usuario", "Pts", "Aciertos", "Exactas"]
    col_widths = [0.3, 1.6, 0.7, 0.8, 0.7]
    hx = left_x + int(0.2*INCH)
    for j, (h, cw) in enumerate(zip(headers, col_widths)):
        add_text_box(slide, hx, table_y, int(cw*INCH), int(0.25*INCH),
                     h, font_size=8, bold=True, color=MUTED)
        hx += int(cw * INCH)
    rank_data = [
        ("#4", "User...", "980", "78%", "12"),
        ("#5", "User...", "920", "72%", "10"),
        ("#6", "User...", "890", "70%", "9"),
    ]
    for ri, row in enumerate(rank_data):
        ry = table_y + int(0.28*INCH) + ri * int(0.25*INCH)
        hx2 = left_x + int(0.2*INCH)
        for j, (val, cw) in enumerate(zip(row, col_widths)):
            add_text_box(slide, hx2, ry, int(cw*INCH), int(0.25*INCH),
                         val, font_size=9, color=SECONDARY if j > 0 else WHITE)
            hx2 += int(cw * INCH)

    # Right panel — Groups
    add_rect(slide, right_x, cursor.y, right_w, panel_h, fill_color=DARK_SURFACE,
             border_color=RGBColor(0x28, 0x30, 0x48))
    add_text_box(slide, right_x + int(0.2*INCH), cursor.y + int(0.15*INCH),
                 int(3*INCH), int(0.4*INCH),
                 "👥  Grupos privados", font_size=18, bold=True, color=WHITE)

    group_items = [
        ("🏟️  Mundial 2026 — Friends", "12 miembros  ·  Top: #3"),
        ("⚽  Los del café", "8 miembros  ·  Top: #1"),
        ("🇦🇷  Albiceleste Squad", "24 miembros  ·  Top: #5"),
    ]
    for gi, (gname, gmeta) in enumerate(group_items):
        gy = cursor.y + int(0.7*INCH) + gi * int(0.7*INCH)
        add_rect(slide, right_x + int(0.2*INCH), gy,
                 right_w - int(0.4*INCH), int(0.6*INCH),
                 fill_color=ELEVATED, border_color=RGBColor(0x28, 0x30, 0x48))
        add_text_box(slide, right_x + int(0.35*INCH), gy + int(0.05*INCH),
                     right_w - int(0.7*INCH), int(0.3*INCH),
                     gname, font_size=13, bold=True, color=WHITE)
        add_text_box(slide, right_x + int(0.35*INCH), gy + int(0.32*INCH),
                     right_w - int(0.7*INCH), int(0.25*INCH),
                     gmeta, font_size=10, color=SECONDARY)

    # Group actions
    actions_y = cursor.y + int(2.85*INCH)
    add_text_box(slide, right_x + int(0.2*INCH), actions_y,
                 right_w - int(0.4*INCH), int(0.25*INCH),
                 "➕  Crear nuevo grupo", font_size=12, bold=True, color=EMERALD)
    add_text_box(slide, right_x + int(0.2*INCH), actions_y + int(0.32*INCH),
                 right_w - int(0.4*INCH), int(0.25*INCH),
                 "🔗  Unirse con código de invitación", font_size=12, bold=True, color=CYAN)

    # Footer rail
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, int(0.65 * INCH),
             fill_color=DARK_SURFACE)
    add_text_box(slide, int(1*INCH), FOOTER_TOP + int(0.12*INCH),
                 SLIDE_W - int(2*INCH), int(0.4*INCH),
                 "06 / 08  ·  Ranking · Grupos privados",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build_slide_7_admin(prs):
    """Admin panel overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MIDNIGHT)
    cursor = Cursor(start_y=int(0.6 * INCH))

    add_rect(slide, 0, 0, int(0.06*INCH), int(0.06*INCH), fill_color=EMERALD)
    add_text_box(slide, int(0.8*INCH), int(0.15*INCH), int(6*INCH), int(0.5*INCH),
                 "Administración", font_size=14, color=MUTED, bold=True)
    add_text_box(slide, int(0.8*INCH), cursor.y, int(10*INCH), int(0.8*INCH),
                 "Panel administrativo", font_size=36, bold=True)
    cursor.advance(int(0.9 * INCH))

    # KPI cards row
    kpis = [
        ("12,480", "Usuarios registrados"),
        ("384", "Grupos creados"),
        ("104", "Partidos programados"),
        ("28.5K", "Predicciones realizadas"),
    ]
    kpi_w = int(2.8 * INCH)
    kpi_gap = int(0.2 * INCH)
    kpi_total = 4 * kpi_w + 3 * kpi_gap
    kpi_start = int((SLIDE_W - kpi_total) / 2)

    for i, (val, label) in enumerate(kpis):
        kx = kpi_start + i * (kpi_w + kpi_gap)
        add_rect(slide, kx, cursor.y, kpi_w, int(1.2*INCH),
                 fill_color=DARK_SURFACE, border_color=RGBColor(0x28, 0x30, 0x48))
        add_text_box(slide, kx, cursor.y + int(0.15*INCH), kpi_w, int(0.5*INCH),
                     val, font_size=28, bold=True, color=EMERALD,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, kx, cursor.y + int(0.65*INCH), kpi_w, int(0.35*INCH),
                     label, font_size=12, color=SECONDARY,
                     alignment=PP_ALIGN.CENTER)

    cursor.advance(int(1.5 * INCH))

    # Admin panels list
    panels = [
        ("⚙️", "Dashboard Admin", "KPIs globales, gráficos de actividad, usuarios activos"),
        ("⚽", "Gestión de partidos", "CRUD de partidos, actualizar estado, registrar resultados"),
        ("👤", "Gestión de usuarios", "Listar, buscar, bloquear/reactivar cuentas"),
        ("📏", "Reglas de puntuación", "Puntos por acierto exacto, ganador, empate"),
    ]

    panel_w = int(5.8 * INCH)
    for i, (icon, pname, pdesc) in enumerate(panels):
        px = int(0.8 * INCH) if i % 2 == 0 else int(7.0 * INCH)
        py = cursor.y + (i // 2) * int(0.8 * INCH)
        add_rect(slide, px, py, panel_w, int(0.65*INCH),
                 fill_color=DARK_SURFACE, border_color=RGBColor(0x28, 0x30, 0x48))
        add_text_box(slide, px + int(0.15*INCH), py + int(0.05*INCH),
                     int(0.4*INCH), int(0.5*INCH),
                     icon, font_size=20)
        add_text_box(slide, px + int(0.6*INCH), py + int(0.05*INCH),
                     int(3*INCH), int(0.3*INCH),
                     pname, font_size=14, bold=True, color=WHITE)
        add_text_box(slide, px + int(0.6*INCH), py + int(0.33*INCH),
                     int(4.5*INCH), int(0.25*INCH),
                     pdesc, font_size=10, color=SECONDARY)

    # Footer rail
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, int(0.65 * INCH),
             fill_color=DARK_SURFACE)
    add_text_box(slide, int(1*INCH), FOOTER_TOP + int(0.12*INCH),
                 SLIDE_W - int(2*INCH), int(0.4*INCH),
                 "07 / 08  ·  Panel administrativo",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build_slide_8_closing(prs):
    """Closing — technical spec & next steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MIDNIGHT)

    # Accent bar
    add_rect(slide, 0, 0, SLIDE_W, int(0.06 * INCH), fill_color=EMERALD)

    # Center content
    content_h = int(3.5 * INCH)
    center_y = int((CONTENT_MAX_Y - content_h) / 2)

    add_text_box(slide, int(1.5*INCH), center_y,
                 SLIDE_W - int(3*INCH), int(0.8*INCH),
                 "Listo para desarrollo", font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, int(2*INCH), center_y + int(1.0*INCH),
                 SLIDE_W - int(4*INCH), int(0.6*INCH),
                 "14 pantallas · Design system completo · Responsive",
                 font_size=20, color=EMERALD, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, int(2.5*INCH), center_y + int(1.7*INCH),
                 SLIDE_W - int(5*INCH), int(0.5*INCH),
                 "Next.js + TypeScript · Componentes reutilizables · Sin apuestas",
                 font_size=14, color=SECONDARY, alignment=PP_ALIGN.CENTER)

    # Tech stack badges
    badges_y = center_y + int(2.5*INCH)
    techs = ["Next.js", "TypeScript", "CSS Tokens", "Responsive", "Sin gambling"]
    badge_w = int(1.3 * INCH)
    badge_gap = int(0.15 * INCH)
    badge_total = 5 * badge_w + 4 * badge_gap
    badge_start = int((SLIDE_W - badge_total) / 2)
    for i, tech in enumerate(techs):
        bx = badge_start + i * (badge_w + badge_gap)
        add_rect(slide, bx, badges_y, badge_w, int(0.35*INCH),
                 fill_color=DARK_SURFACE, border_color=RGBColor(0x28, 0x30, 0x48))
        add_text_box(slide, bx, badges_y + int(0.03*INCH),
                     badge_w, int(0.3*INCH),
                     tech, font_size=9, color=CYAN, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Footer rail
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, int(0.65 * INCH),
             fill_color=DARK_SURFACE)
    add_text_box(slide, int(1*INCH), FOOTER_TOP + int(0.12*INCH),
                 SLIDE_W - int(2*INCH), int(0.4*INCH),
                 "08 / 08  ·  Quiniela Mundial 2026 · Próximos pasos",
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


# ── Main ────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1_cover(prs)
    build_slide_2_design_system(prs)
    build_slide_3_public_screens(prs)
    build_slide_4_dashboard(prs)
    build_slide_5_matches_predictions(prs)
    build_slide_6_leaderboard_groups(prs)
    build_slide_7_admin(prs)
    build_slide_8_closing(prs)

    output = "quiniela-mundial-2026.pptx"
    prs.save(output)
    print("[OK] Saved: " + output)
    print("   Slides: " + str(len(prs.slides)))
    print("   Rail: CONTENT_MAX_Y={:.2f}in, FOOTER_TOP={:.2f}in".format(CONTENT_MAX_Y/INCH, FOOTER_TOP/INCH))
    print("   Footer-rail discipline applied to all slides")
    print("   a:latin = '{}', a:ea = '{}' on all runs".format(LATIN_FACE, EA_FACE))

if __name__ == "__main__":
    main()
